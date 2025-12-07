"""
Advanced PDF to PowerPoint Converter
Uses AI and sophisticated analysis to convert PDFs to editable PowerPoint presentations
with preserved fonts, styles, and live charts.
"""

import os
import io
import json
import tempfile
from typing import List, Dict, Any, Tuple, Optional
from pathlib import Path
from dataclasses import dataclass
from collections import defaultdict

import pdfplumber
import fitz  # PyMuPDF
from PIL import Image
import numpy as np
import cv2
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
from matplotlib import font_manager


@dataclass
class TextElement:
    """Represents a text element from the PDF"""
    text: str
    x: float
    y: float
    width: float
    height: float
    font_name: str
    font_size: float
    color: Tuple[int, int, int]
    bold: bool = False
    italic: bool = False
    alignment: str = 'left'


@dataclass
class ImageElement:
    """Represents an image element from the PDF"""
    image: Image.Image
    x: float
    y: float
    width: float
    height: float
    name: str


@dataclass
class ShapeElement:
    """Represents a shape/vector element from the PDF"""
    shape_type: str
    x: float
    y: float
    width: float
    height: float
    fill_color: Optional[Tuple[int, int, int]]
    stroke_color: Optional[Tuple[int, int, int]]
    stroke_width: float
    points: List[Tuple[float, float]] = None


@dataclass
class ChartElement:
    """Represents a detected chart/graph from the PDF"""
    chart_type: str  # 'bar', 'line', 'pie', 'scatter', etc.
    x: float
    y: float
    width: float
    height: float
    title: str
    data: Dict[str, Any]
    image: Image.Image
    categories: List[str]
    series: List[Dict[str, Any]]


@dataclass
class PageStructure:
    """Represents the complete structure of a PDF page"""
    page_number: int
    width: float
    height: float
    text_elements: List[TextElement]
    image_elements: List[ImageElement]
    shape_elements: List[ShapeElement]
    chart_elements: List[ChartElement]
    background_color: Tuple[int, int, int]


class PDFAnalyzer:
    """Analyzes PDF structure and extracts elements"""
    
    def __init__(self, pdf_path: str):
        self.pdf_path = pdf_path
        self.plumber_pdf = pdfplumber.open(pdf_path)
        self.fitz_pdf = fitz.open(pdf_path)
    
    def analyze_page(self, page_num: int) -> PageStructure:
        """Analyze a single page and extract all elements"""
        plumber_page = self.plumber_pdf.pages[page_num]
        fitz_page = self.fitz_pdf[page_num]
        
        # Get page dimensions
        width = plumber_page.width
        height = plumber_page.height
        
        # Extract different elements
        text_elements = self._extract_text_elements(plumber_page, fitz_page)
        image_elements = self._extract_images(fitz_page)
        shape_elements = self._extract_shapes(fitz_page)
        
        # Detect background color
        bg_color = self._detect_background_color(fitz_page)
        
        return PageStructure(
            page_number=page_num,
            width=width,
            height=height,
            text_elements=text_elements,
            image_elements=image_elements,
            shape_elements=shape_elements,
            chart_elements=[],  # Will be populated by AI analysis
            background_color=bg_color
        )
    
    def _extract_text_elements(self, plumber_page, fitz_page) -> List[TextElement]:
        """Extract text with detailed formatting information"""
        text_elements = []
        
        # Use pdfplumber for detailed text extraction
        chars = plumber_page.chars
        
        # Group characters into words and lines
        current_line = []
        current_y = None
        tolerance = 2  # pixels
        
        for char in chars:
            y = char['y0']
            
            # Check if this is a new line
            if current_y is None:
                current_y = y
            elif abs(y - current_y) > tolerance:
                # Process current line
                if current_line:
                    text_elements.extend(self._process_line(current_line))
                current_line = [char]
                current_y = y
            else:
                current_line.append(char)
        
        # Process last line
        if current_line:
            text_elements.extend(self._process_line(current_line))
        
        return text_elements
    
    def _process_line(self, chars: List[Dict]) -> List[TextElement]:
        """Process a line of characters into text elements"""
        if not chars:
            return []
        
        elements = []
        current_word = []
        current_font = None
        
        for i, char in enumerate(chars):
            font_key = (char.get('fontname', ''), char.get('size', 12))
            
            # Check if font changed or space detected
            next_char = chars[i + 1] if i + 1 < len(chars) else None
            is_space = next_char and (next_char['x0'] - char['x1'] > char['width'])
            
            if current_font is None:
                current_font = font_key
            
            if font_key != current_font or is_space or i == len(chars) - 1:
                # Create text element from current word
                if current_word or (i == len(chars) - 1 and char['text'].strip()):
                    if i == len(chars) - 1 and font_key == current_font:
                        current_word.append(char)
                    
                    if current_word:
                        text = ''.join([c['text'] for c in current_word])
                        x0 = min([c['x0'] for c in current_word])
                        y0 = min([c['y0'] for c in current_word])
                        x1 = max([c['x1'] for c in current_word])
                        y1 = max([c['y1'] for c in current_word])
                        
                        font_name = current_word[0].get('fontname', 'Arial')
                        font_size = current_word[0].get('size', 12)
                        
                        # Detect bold/italic from font name
                        bold = 'bold' in font_name.lower()
                        italic = 'italic' in font_name.lower() or 'oblique' in font_name.lower()
                        
                        # Extract color (default to black)
                        color = self._get_text_color(current_word[0])
                        
                        elements.append(TextElement(
                            text=text,
                            x=x0,
                            y=y0,
                            width=x1 - x0,
                            height=y1 - y0,
                            font_name=self._normalize_font_name(font_name),
                            font_size=font_size,
                            color=color,
                            bold=bold,
                            italic=italic
                        ))
                
                # Start new word
                if not is_space and font_key != current_font:
                    current_word = [char]
                    current_font = font_key
                else:
                    current_word = []
                    current_font = None
            else:
                current_word.append(char)
        
        return elements
    
    def _get_text_color(self, char: Dict) -> Tuple[int, int, int]:
        """Extract text color from character"""
        # pdfplumber doesn't always provide color info
        # Default to black
        return (0, 0, 0)
    
    def _normalize_font_name(self, font_name: str) -> str:
        """Normalize PDF font names to standard font names"""
        font_name = font_name.lower()
        
        # Common font mappings
        if 'arial' in font_name:
            return 'Arial'
        elif 'helvetica' in font_name:
            return 'Helvetica'
        elif 'times' in font_name:
            return 'Times New Roman'
        elif 'courier' in font_name:
            return 'Courier New'
        elif 'calibri' in font_name:
            return 'Calibri'
        elif 'verdana' in font_name:
            return 'Verdana'
        elif 'georgia' in font_name:
            return 'Georgia'
        elif 'comic' in font_name:
            return 'Comic Sans MS'
        else:
            # Return capitalized version
            return font_name.split('+')[-1].capitalize()
    
    def _extract_images(self, fitz_page) -> List[ImageElement]:
        """Extract images from the page"""
        image_elements = []
        image_list = fitz_page.get_images()
        
        for img_index, img in enumerate(image_list):
            xref = img[0]
            
            try:
                # Get image bbox
                img_rects = fitz_page.get_image_rects(xref)
                if not img_rects:
                    continue
                
                rect = img_rects[0]
                
                # Extract image
                base_image = self.fitz_pdf.extract_image(xref)
                image_bytes = base_image["image"]
                
                # Convert to PIL Image
                pil_image = Image.open(io.BytesIO(image_bytes))
                
                image_elements.append(ImageElement(
                    image=pil_image,
                    x=rect.x0,
                    y=rect.y0,
                    width=rect.width,
                    height=rect.height,
                    name=f"image_{img_index}"
                ))
            except Exception as e:
                print(f"Warning: Could not extract image {img_index}: {e}")
                continue
        
        return image_elements
    
    def _extract_shapes(self, fitz_page) -> List[ShapeElement]:
        """Extract vector shapes from the page"""
        shape_elements = []
        
        # Get drawing commands
        drawings = fitz_page.get_drawings()
        
        for drawing in drawings:
            rect = drawing.get('rect')
            if not rect:
                continue
            
            fill_color = drawing.get('fill')
            stroke_color = drawing.get('color')
            width = drawing.get('width', 1)
            
            # Convert colors to RGB tuples
            fill_rgb = self._convert_color(fill_color) if fill_color else None
            stroke_rgb = self._convert_color(stroke_color) if stroke_color else None
            
            # Determine shape type
            shape_type = drawing.get('type', 'path')
            
            shape_elements.append(ShapeElement(
                shape_type=shape_type,
                x=rect.x0,
                y=rect.y0,
                width=rect.width,
                height=rect.height,
                fill_color=fill_rgb,
                stroke_color=stroke_rgb,
                stroke_width=width
            ))
        
        return shape_elements
    
    def _convert_color(self, color) -> Tuple[int, int, int]:
        """Convert PDF color to RGB tuple"""
        if isinstance(color, (list, tuple)):
            if len(color) == 3:
                return tuple(int(c * 255) for c in color)
            elif len(color) == 1:
                # Grayscale
                return tuple([int(color[0] * 255)] * 3)
        return (0, 0, 0)
    
    def _detect_background_color(self, fitz_page) -> Tuple[int, int, int]:
        """Detect the background color of the page"""
        # Render page to image
        pix = fitz_page.get_pixmap(alpha=False)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        
        # Get most common color from corners
        img_array = np.array(img)
        
        # Sample corners
        corner_size = 50
        corners = [
            img_array[0:corner_size, 0:corner_size],
            img_array[0:corner_size, -corner_size:],
            img_array[-corner_size:, 0:corner_size],
            img_array[-corner_size:, -corner_size:]
        ]
        
        # Get average color from corners
        avg_colors = []
        for corner in corners:
            avg_color = corner.reshape(-1, 3).mean(axis=0)
            avg_colors.append(avg_color)
        
        # Return average of all corners
        final_color = np.mean(avg_colors, axis=0).astype(int)
        return tuple(final_color)
    
    def close(self):
        """Close PDF files"""
        self.plumber_pdf.close()
        self.fitz_pdf.close()


class ChartDetector:
    """Detects and analyzes charts in PDF pages using AI and image analysis"""
    
    def __init__(self, ai_utils_module):
        self.ai_utils = ai_utils_module
    
    def detect_charts(self, page_structure: PageStructure, page_image: Image.Image) -> List[ChartElement]:
        """Detect charts in the page using AI analysis"""
        charts = []
        
        # First, use image analysis to find potential chart regions
        potential_regions = self._find_chart_regions(page_image, page_structure)
        
        # Then use AI to analyze each region
        for region in potential_regions:
            try:
                chart = self._analyze_chart_with_ai(region, page_image)
                if chart:
                    charts.append(chart)
            except Exception as e:
                print(f"Warning: Could not analyze chart region: {e}")
        
        return charts
    
    def _find_chart_regions(self, page_image: Image.Image, page_structure: PageStructure) -> List[Dict]:
        """Find potential chart regions using image analysis and structure"""
        regions = []
        
        # Convert to numpy array
        img_array = np.array(page_image)
        
        # Look for image elements that might be charts
        for img_elem in page_structure.image_elements:
            # Charts are typically larger images
            if img_elem.width > 100 and img_elem.height > 100:
                regions.append({
                    'x': img_elem.x,
                    'y': img_elem.y,
                    'width': img_elem.width,
                    'height': img_elem.height,
                    'image': img_elem.image
                })
        
        # Also look for grouped shapes that might form a chart
        # (This is a simplified heuristic)
        shape_groups = self._group_nearby_shapes(page_structure.shape_elements)
        for group in shape_groups:
            if len(group) > 5:  # Charts typically have multiple shapes
                x = min([s.x for s in group])
                y = min([s.y for s in group])
                x2 = max([s.x + s.width for s in group])
                y2 = max([s.y + s.height for s in group])
                width = x2 - x
                height = y2 - y
                
                if width > 100 and height > 100:
                    # Crop the region from page image
                    # Convert coordinates
                    scale_x = img_array.shape[1] / page_structure.width
                    scale_y = img_array.shape[0] / page_structure.height
                    
                    x_pixel = int(x * scale_x)
                    y_pixel = int(y * scale_y)
                    w_pixel = int(width * scale_x)
                    h_pixel = int(height * scale_y)
                    
                    cropped = img_array[y_pixel:y_pixel+h_pixel, x_pixel:x_pixel+w_pixel]
                    cropped_img = Image.fromarray(cropped)
                    
                    regions.append({
                        'x': x,
                        'y': y,
                        'width': width,
                        'height': height,
                        'image': cropped_img
                    })
        
        return regions
    
    def _group_nearby_shapes(self, shapes: List[ShapeElement], distance_threshold: float = 50) -> List[List[ShapeElement]]:
        """Group shapes that are close together"""
        if not shapes:
            return []
        
        groups = []
        used = set()
        
        for i, shape in enumerate(shapes):
            if i in used:
                continue
            
            group = [shape]
            used.add(i)
            
            # Find nearby shapes
            for j, other_shape in enumerate(shapes):
                if j in used:
                    continue
                
                # Calculate distance
                dist = self._shape_distance(shape, other_shape)
                if dist < distance_threshold:
                    group.append(other_shape)
                    used.add(j)
            
            groups.append(group)
        
        return groups
    
    def _shape_distance(self, shape1: ShapeElement, shape2: ShapeElement) -> float:
        """Calculate distance between two shapes"""
        center1_x = shape1.x + shape1.width / 2
        center1_y = shape1.y + shape1.height / 2
        center2_x = shape2.x + shape2.width / 2
        center2_y = shape2.y + shape2.height / 2
        
        return ((center1_x - center2_x)**2 + (center1_y - center2_y)**2)**0.5
    
    def _analyze_chart_with_ai(self, region: Dict, page_image: Image.Image) -> Optional[ChartElement]:
        """Use AI to analyze a potential chart region"""
        try:
            # Use AI to analyze the chart image
            chart_data = self.ai_utils.analyze_chart_image(
                chart_image=region['image'],
                context="This is a chart extracted from a PDF page."
            )
            
            # Check if we got valid chart data
            if not chart_data.get('categories') or not chart_data.get('series'):
                return None
            
            # Create ChartElement
            return ChartElement(
                chart_type=chart_data.get('chart_type', 'column'),
                x=region['x'],
                y=region['y'],
                width=region['width'],
                height=region['height'],
                title=chart_data.get('title', ''),
                data=chart_data,
                image=region['image'],
                categories=chart_data.get('categories', []),
                series=chart_data.get('series', [])
            )
        except Exception as e:
            print(f"Warning: Failed to analyze chart with AI: {e}")
            return None


class PowerPointBuilder:
    """Builds PowerPoint presentations from analyzed PDF structure"""
    
    def __init__(self):
        self.prs = Presentation()
        # Set standard slide dimensions (16:9)
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
    
    def add_page(self, page_structure: PageStructure, ai_enhanced_data: Optional[Dict] = None):
        """Add a slide from page structure"""
        # Use blank layout
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # Set background
        self._set_background(slide, page_structure.background_color)
        
        # Calculate scaling factors (PDF points to PowerPoint EMUs)
        # PowerPoint slide dimensions are in EMUs
        # PDF dimensions are in points (72 points = 1 inch)
        # We need simple ratio scaling, not EMU conversion
        scale_x = float(self.prs.slide_width) / float(page_structure.width)
        scale_y = float(self.prs.slide_height) / float(page_structure.height)
        
        # Add shapes first (they go to the back)
        for shape_elem in page_structure.shape_elements:
            self._add_shape(slide, shape_elem, scale_x, scale_y)
        
        # Add images
        for img_elem in page_structure.image_elements:
            self._add_image(slide, img_elem, scale_x, scale_y)
        
        # Add charts (if any)
        for chart_elem in page_structure.chart_elements:
            self._add_chart(slide, chart_elem, scale_x, scale_y)
        
        # Add text (on top)
        for text_elem in page_structure.text_elements:
            self._add_text(slide, text_elem, scale_x, scale_y)
    
    def _set_background(self, slide, color: Tuple[int, int, int]):
        """Set slide background color"""
        background = slide.background
        fill = background.fill
        fill.solid()
        # Ensure RGB values are integers
        r = int(color[0])
        g = int(color[1])
        b = int(color[2])
        fill.fore_color.rgb = RGBColor(r, g, b)
    
    def _add_shape(self, slide, shape_elem: ShapeElement, scale_x: float, scale_y: float):
        """Add a shape to the slide"""
        left = int(shape_elem.x * scale_x)
        top = int(shape_elem.y * scale_y)
        width = int(shape_elem.width * scale_x)
        height = int(shape_elem.height * scale_y)
        
        # Add rectangle (simplified - could be extended for other shapes)
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            left, top, width, height
        )
        
        # Set fill color
        if shape_elem.fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*shape_elem.fill_color)
        else:
            shape.fill.background()
        
        # Set line color
        if shape_elem.stroke_color:
            shape.line.color.rgb = RGBColor(*shape_elem.stroke_color)
            shape.line.width = Pt(shape_elem.stroke_width)
        else:
            shape.line.fill.background()
    
    def _add_image(self, slide, img_elem: ImageElement, scale_x: float, scale_y: float):
        """Add an image to the slide"""
        left = int(img_elem.x * scale_x)
        top = int(img_elem.y * scale_y)
        width = int(img_elem.width * scale_x)
        height = int(img_elem.height * scale_y)
        
        # Save image to temporary file
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            img_elem.image.save(tmp.name, 'PNG')
            slide.shapes.add_picture(tmp.name, left, top, width, height)
            os.unlink(tmp.name)
    
    def _add_text(self, slide, text_elem: TextElement, scale_x: float, scale_y: float):
        """Add a text box to the slide"""
        left = int(text_elem.x * scale_x)
        top = int(text_elem.y * scale_y)
        width = int(text_elem.width * scale_x)
        height = int(text_elem.height * scale_y)
        
        # Add text box
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = False
        
        # Add text
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text_elem.text
        
        # Set font
        font = run.font
        font.name = text_elem.font_name
        
        # Calculate font size - scale_x and scale_y are already ratios
        # PowerPoint expects EMUs for dimensions but Points for font size
        # Since both PDF and PPT use points for fonts, we can use the font size directly
        # with minimal scaling for aspect ratio differences
        font_scale = min(scale_x / (self.prs.slide_width / Inches(10)),
                        scale_y / (self.prs.slide_height / Inches(7.5)))
        scaled_font_size = text_elem.font_size * font_scale
        
        # Ensure font size is within valid range (1-409 points)
        scaled_font_size = max(1, min(409, scaled_font_size))
        font.size = Pt(scaled_font_size)
        
        font.bold = text_elem.bold
        font.italic = text_elem.italic
        font.color.rgb = RGBColor(int(text_elem.color[0]), int(text_elem.color[1]), int(text_elem.color[2]))
        
        # Set alignment
        if text_elem.alignment == 'center':
            p.alignment = PP_ALIGN.CENTER
        elif text_elem.alignment == 'right':
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT
    
    def _add_chart(self, slide, chart_elem: ChartElement, scale_x: float, scale_y: float):
        """Add an editable chart to the slide"""
        left = int(chart_elem.x * scale_x)
        top = int(chart_elem.y * scale_y)
        width = int(chart_elem.width * scale_x)
        height = int(chart_elem.height * scale_y)
        
        # Determine chart type
        chart_type = self._get_pptx_chart_type(chart_elem.chart_type)
        
        # Create chart data
        chart_data = CategoryChartData()
        chart_data.categories = chart_elem.categories
        
        for series in chart_elem.series:
            chart_data.add_series(series['name'], series['values'])
        
        # Add chart
        graphic_frame = slide.shapes.add_chart(
            chart_type, left, top, width, height, chart_data
        )
        
        chart = graphic_frame.chart
        
        # Set chart title
        if chart_elem.title:
            chart.has_title = True
            chart.chart_title.text_frame.text = chart_elem.title
        
        # Set legend
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    
    def _get_pptx_chart_type(self, chart_type_str: str):
        """Convert chart type string to PowerPoint chart type"""
        chart_type_map = {
            'bar': XL_CHART_TYPE.BAR_CLUSTERED,
            'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'line': XL_CHART_TYPE.LINE,
            'pie': XL_CHART_TYPE.PIE,
            'scatter': XL_CHART_TYPE.XY_SCATTER,
            'area': XL_CHART_TYPE.AREA,
        }
        return chart_type_map.get(chart_type_str.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)
    
    def save(self, output_path: str):
        """Save the presentation"""
        self.prs.save(output_path)


def convert_pdf_to_pptx(
    pdf_path: str,
    output_path: str,
    use_ai: bool = True,
    ai_utils_module = None,
    progress_callback = None
) -> str:
    """
    Convert a PDF to PowerPoint presentation
    
    Args:
        pdf_path: Path to input PDF
        output_path: Path to output PPTX
        use_ai: Whether to use AI for enhanced conversion
        ai_utils_module: AI utilities module for enhanced conversion
        progress_callback: Callback function for progress updates
    
    Returns:
        Path to the created PPTX file
    """
    # Initialize components
    analyzer = PDFAnalyzer(pdf_path)
    builder = PowerPointBuilder()
    
    if use_ai and ai_utils_module:
        chart_detector = ChartDetector(ai_utils_module)
    else:
        chart_detector = None
    
    try:
        # Get total pages
        total_pages = len(analyzer.plumber_pdf.pages)
        
        if progress_callback:
            progress_callback(f"Converting {total_pages} pages...")
        
        # Process each page
        for page_num in range(total_pages):
            if progress_callback:
                progress_callback(f"Processing page {page_num + 1}/{total_pages}...")
            
            # Analyze page structure
            page_structure = analyzer.analyze_page(page_num)
            
            # Detect charts with AI if enabled
            if chart_detector:
                try:
                    # Render page to image for chart detection
                    page_image = Image.open(io.BytesIO(
                        analyzer.fitz_pdf[page_num].get_pixmap().tobytes("png")
                    ))
                    page_structure.chart_elements = chart_detector.detect_charts(
                        page_structure, page_image
                    )
                except Exception as e:
                    print(f"Warning: Chart detection failed for page {page_num + 1}: {e}")
            
            # Add to presentation
            builder.add_page(page_structure)
        
        # Save presentation
        if progress_callback:
            progress_callback("Saving PowerPoint presentation...")
        
        builder.save(output_path)
        
        if progress_callback:
            progress_callback(f"Conversion complete! Saved to {output_path}")
        
        return output_path
    
    finally:
        analyzer.close()
