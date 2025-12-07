"""
PowerPoint utilities for creating and manipulating PPTX files.
Provides functions to create slides, add text boxes, charts, images, and more.
"""

import os
import io
import re
import tempfile
from typing import List, Dict, Tuple, Optional, Any, Union
from dataclasses import dataclass, field
from enum import Enum
from PIL import Image
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml


# ============================================================================
# Data Classes for Slide Elements
# ============================================================================

class ChartType(Enum):
    """Supported chart types for conversion"""
    BAR = "bar"
    COLUMN = "column"
    LINE = "line"
    PIE = "pie"
    DONUT = "donut"
    AREA = "area"
    SCATTER = "scatter"
    BUBBLE = "bubble"
    STACKED_BAR = "stacked_bar"
    STACKED_COLUMN = "stacked_column"
    CLUSTERED_BAR = "clustered_bar"
    CLUSTERED_COLUMN = "clustered_column"
    WATERFALL = "waterfall"
    COMBO = "combo"


@dataclass
class TextStyle:
    """Text styling properties"""
    font_name: str = "Calibri"
    font_size: int = 12  # in points
    bold: bool = False
    italic: bool = False
    underline: bool = False
    color: str = "#000000"  # hex color
    alignment: str = "left"  # left, center, right, justify
    vertical_alignment: str = "top"  # top, middle, bottom
    line_spacing: float = 1.0
    space_before: float = 0  # points
    space_after: float = 0  # points


@dataclass
class TextRun:
    """A run of text with consistent styling"""
    text: str
    style: TextStyle = field(default_factory=TextStyle)


@dataclass
class Paragraph:
    """A paragraph containing multiple text runs"""
    runs: List[TextRun] = field(default_factory=list)
    alignment: str = "left"
    bullet: bool = False
    bullet_char: str = "•"
    level: int = 0  # indentation level


@dataclass
class TextBox:
    """A text box element"""
    paragraphs: List[Paragraph]
    x: float  # inches from left
    y: float  # inches from top
    width: float  # inches
    height: float  # inches
    background_color: Optional[str] = None  # hex color
    border_color: Optional[str] = None
    border_width: float = 0  # points
    rotation: float = 0  # degrees
    margin_left: float = 0.1  # inches
    margin_right: float = 0.1
    margin_top: float = 0.05
    margin_bottom: float = 0.05


@dataclass
class ChartSeries:
    """A data series for charts"""
    name: str
    values: List[float]
    color: Optional[str] = None  # hex color
    x_values: Optional[List[float]] = None  # for scatter/bubble charts
    sizes: Optional[List[float]] = None  # for bubble charts


@dataclass
class ChartData:
    """Data for creating charts"""
    chart_type: ChartType
    title: Optional[str] = None
    categories: List[str] = field(default_factory=list)
    series: List[ChartSeries] = field(default_factory=list)
    x_axis_title: Optional[str] = None
    y_axis_title: Optional[str] = None
    show_legend: bool = True
    legend_position: str = "bottom"  # top, bottom, left, right
    data_labels: bool = False
    colors: List[str] = field(default_factory=list)  # custom color palette


@dataclass
class ChartElement:
    """A chart element on a slide"""
    data: ChartData
    x: float  # inches from left
    y: float  # inches from top
    width: float  # inches
    height: float  # inches


@dataclass
class ImageElement:
    """An image element on a slide"""
    image_data: bytes  # raw image bytes
    x: float  # inches from left
    y: float  # inches from top
    width: float  # inches
    height: float  # inches
    alt_text: str = ""
    rotation: float = 0  # degrees
    crop_left: float = 0
    crop_right: float = 0
    crop_top: float = 0
    crop_bottom: float = 0


@dataclass
class ShapeElement:
    """A shape element (rectangle, circle, arrow, etc.)"""
    shape_type: str  # rectangle, oval, triangle, arrow, etc.
    x: float
    y: float
    width: float
    height: float
    fill_color: Optional[str] = None
    border_color: Optional[str] = None
    border_width: float = 1  # points
    rotation: float = 0
    text: Optional[str] = None
    text_style: Optional[TextStyle] = None


@dataclass
class TableCell:
    """A cell in a table"""
    text: str
    style: TextStyle = field(default_factory=TextStyle)
    background_color: Optional[str] = None
    colspan: int = 1
    rowspan: int = 1


@dataclass
class TableElement:
    """A table element on a slide"""
    rows: List[List[TableCell]]
    x: float
    y: float
    width: float
    height: float
    header_row: bool = True
    header_color: Optional[str] = None
    border_color: str = "#000000"
    alternating_row_color: Optional[str] = None


@dataclass
class SlideLayout:
    """Slide layout information"""
    width: float = 13.333  # inches (16:9 default)
    height: float = 7.5  # inches
    background_color: Optional[str] = None
    background_image: Optional[bytes] = None


@dataclass
class SlideContent:
    """Complete content for a single slide"""
    layout: SlideLayout = field(default_factory=SlideLayout)
    text_boxes: List[TextBox] = field(default_factory=list)
    charts: List[ChartElement] = field(default_factory=list)
    images: List[ImageElement] = field(default_factory=list)
    shapes: List[ShapeElement] = field(default_factory=list)
    tables: List[TableElement] = field(default_factory=list)
    notes: str = ""


# ============================================================================
# Color Utilities
# ============================================================================

def hex_to_rgb(hex_color: str) -> Tuple[int, int, int]:
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    if len(hex_color) == 3:
        hex_color = ''.join([c*2 for c in hex_color])
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def rgb_to_hex(r: int, g: int, b: int) -> str:
    """Convert RGB values to hex color"""
    return f"#{r:02x}{g:02x}{b:02x}"


def get_rgb_color(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor object"""
    r, g, b = hex_to_rgb(hex_color)
    return RGBColor(r, g, b)


# ============================================================================
# PowerPoint Creation Functions
# ============================================================================

def create_presentation(
    slides_content: List[SlideContent],
    slide_width: float = 13.333,
    slide_height: float = 7.5
) -> Presentation:
    """
    Create a PowerPoint presentation from slide content.
    
    Args:
        slides_content: List of SlideContent objects
        slide_width: Width of slides in inches (default 16:9 ratio)
        slide_height: Height of slides in inches
    
    Returns:
        Presentation object
    """
    prs = Presentation()
    
    # Set slide dimensions
    prs.slide_width = Inches(slide_width)
    prs.slide_height = Inches(slide_height)
    
    # Get blank layout
    blank_layout = prs.slide_layouts[6]  # Blank layout
    
    for slide_content in slides_content:
        slide = prs.slides.add_slide(blank_layout)
        
        # Set background if specified
        if slide_content.layout.background_color:
            _set_slide_background_color(slide, slide_content.layout.background_color)
        elif slide_content.layout.background_image:
            _set_slide_background_image(slide, slide_content.layout.background_image)
        
        # Add text boxes
        for text_box in slide_content.text_boxes:
            _add_text_box(slide, text_box)
        
        # Add charts
        for chart_elem in slide_content.charts:
            _add_chart(slide, chart_elem)
        
        # Add images
        for image_elem in slide_content.images:
            _add_image(slide, image_elem)
        
        # Add shapes
        for shape_elem in slide_content.shapes:
            _add_shape(slide, shape_elem)
        
        # Add tables
        for table_elem in slide_content.tables:
            _add_table(slide, table_elem)
        
        # Add notes
        if slide_content.notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_content.notes
    
    return prs


def save_presentation(prs: Presentation, output_path: str):
    """Save presentation to file"""
    prs.save(output_path)


def _set_slide_background_color(slide, hex_color: str):
    """Set solid background color for a slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = get_rgb_color(hex_color)


def _set_slide_background_image(slide, image_data: bytes):
    """Set background image for a slide"""
    # Save image temporarily
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as f:
        f.write(image_data)
        temp_path = f.name
    
    try:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Add image as background
        left = Inches(0)
        top = Inches(0)
        slide.shapes.add_picture(
            temp_path, left, top,
            width=slide.shapes._spTree.getparent().getparent().get('cx'),
            height=slide.shapes._spTree.getparent().getparent().get('cy')
        )
    finally:
        os.unlink(temp_path)


def _get_alignment(align_str: str) -> PP_ALIGN:
    """Convert alignment string to PP_ALIGN enum"""
    alignments = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT,
        'justify': PP_ALIGN.JUSTIFY,
    }
    return alignments.get(align_str.lower(), PP_ALIGN.LEFT)


def _get_vertical_anchor(valign_str: str) -> MSO_ANCHOR:
    """Convert vertical alignment string to MSO_ANCHOR enum"""
    anchors = {
        'top': MSO_ANCHOR.TOP,
        'middle': MSO_ANCHOR.MIDDLE,
        'bottom': MSO_ANCHOR.BOTTOM,
    }
    return anchors.get(valign_str.lower(), MSO_ANCHOR.TOP)


def _add_text_box(slide, text_box: TextBox):
    """Add a text box to the slide"""
    shape = slide.shapes.add_textbox(
        Inches(text_box.x),
        Inches(text_box.y),
        Inches(text_box.width),
        Inches(text_box.height)
    )
    
    # Set rotation if specified
    if text_box.rotation:
        shape.rotation = text_box.rotation
    
    # Set margins
    text_frame = shape.text_frame
    text_frame.margin_left = Inches(text_box.margin_left)
    text_frame.margin_right = Inches(text_box.margin_right)
    text_frame.margin_top = Inches(text_box.margin_top)
    text_frame.margin_bottom = Inches(text_box.margin_bottom)
    text_frame.word_wrap = True
    
    # Set background color
    if text_box.background_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = get_rgb_color(text_box.background_color)
    else:
        shape.fill.background()  # Transparent
    
    # Set border
    if text_box.border_color:
        shape.line.color.rgb = get_rgb_color(text_box.border_color)
        shape.line.width = Pt(text_box.border_width)
    else:
        shape.line.fill.background()
    
    # Clear default paragraph and add content
    first = True
    for para_data in text_box.paragraphs:
        if first:
            p = text_frame.paragraphs[0]
            first = False
        else:
            p = text_frame.add_paragraph()
        
        # Set paragraph alignment
        p.alignment = _get_alignment(para_data.alignment)
        p.level = para_data.level
        
        # Add bullet if specified
        if para_data.bullet:
            p.level = para_data.level
            # Bullet formatting is typically controlled by the layout
        
        # Add text runs
        first_run = True
        for run_data in para_data.runs:
            if first_run:
                run = p.runs[0] if p.runs else p.add_run()
                first_run = False
            else:
                run = p.add_run()
            
            run.text = run_data.text
            
            # Apply styling
            font = run.font
            font.name = run_data.style.font_name
            font.size = Pt(run_data.style.font_size)
            font.bold = run_data.style.bold
            font.italic = run_data.style.italic
            font.underline = run_data.style.underline
            font.color.rgb = get_rgb_color(run_data.style.color)


def _get_chart_type(chart_type: ChartType) -> XL_CHART_TYPE:
    """Convert ChartType enum to XL_CHART_TYPE"""
    type_map = {
        ChartType.BAR: XL_CHART_TYPE.BAR_CLUSTERED,
        ChartType.COLUMN: XL_CHART_TYPE.COLUMN_CLUSTERED,
        ChartType.LINE: XL_CHART_TYPE.LINE,
        ChartType.PIE: XL_CHART_TYPE.PIE,
        ChartType.DONUT: XL_CHART_TYPE.DOUGHNUT,
        ChartType.AREA: XL_CHART_TYPE.AREA,
        ChartType.SCATTER: XL_CHART_TYPE.XY_SCATTER,
        ChartType.BUBBLE: XL_CHART_TYPE.BUBBLE,
        ChartType.STACKED_BAR: XL_CHART_TYPE.BAR_STACKED,
        ChartType.STACKED_COLUMN: XL_CHART_TYPE.COLUMN_STACKED,
        ChartType.CLUSTERED_BAR: XL_CHART_TYPE.BAR_CLUSTERED,
        ChartType.CLUSTERED_COLUMN: XL_CHART_TYPE.COLUMN_CLUSTERED,
        ChartType.WATERFALL: XL_CHART_TYPE.COLUMN_CLUSTERED,  # Waterfall uses special formatting
        ChartType.COMBO: XL_CHART_TYPE.COLUMN_CLUSTERED,  # Combo charts need special handling
    }
    return type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)


def _get_legend_position(pos_str: str) -> XL_LEGEND_POSITION:
    """Convert legend position string to XL_LEGEND_POSITION"""
    positions = {
        'top': XL_LEGEND_POSITION.TOP,
        'bottom': XL_LEGEND_POSITION.BOTTOM,
        'left': XL_LEGEND_POSITION.LEFT,
        'right': XL_LEGEND_POSITION.RIGHT,
        'corner': XL_LEGEND_POSITION.CORNER,
    }
    return positions.get(pos_str.lower(), XL_LEGEND_POSITION.BOTTOM)


def _add_chart(slide, chart_elem: ChartElement):
    """Add a chart to the slide"""
    chart_data = chart_elem.data
    xl_chart_type = _get_chart_type(chart_data.chart_type)
    
    # Handle different chart types
    if chart_data.chart_type in [ChartType.SCATTER, ChartType.BUBBLE]:
        data = _create_xy_chart_data(chart_data)
    else:
        data = _create_category_chart_data(chart_data)
    
    # Create the chart
    x, y = Inches(chart_elem.x), Inches(chart_elem.y)
    cx, cy = Inches(chart_elem.width), Inches(chart_elem.height)
    
    graphic_frame = slide.shapes.add_chart(xl_chart_type, x, y, cx, cy, data)
    chart = graphic_frame.chart
    
    # Set chart title
    if chart_data.title:
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_data.title
    else:
        chart.has_title = False
    
    # Set legend
    if chart_data.show_legend:
        chart.has_legend = True
        chart.legend.position = _get_legend_position(chart_data.legend_position)
        chart.legend.include_in_layout = False
    else:
        chart.has_legend = False
    
    # Set axis titles if applicable
    if chart_data.chart_type not in [ChartType.PIE, ChartType.DONUT]:
        try:
            if chart_data.x_axis_title:
                chart.category_axis.has_title = True
                chart.category_axis.axis_title.text_frame.text = chart_data.x_axis_title
            if chart_data.y_axis_title:
                chart.value_axis.has_title = True
                chart.value_axis.axis_title.text_frame.text = chart_data.y_axis_title
        except:
            pass  # Some chart types don't support axis titles
    
    # Apply custom colors to series
    if chart_data.colors:
        _apply_chart_colors(chart, chart_data.colors)
    else:
        # Apply colors from individual series if specified
        for i, series in enumerate(chart_data.series):
            if series.color:
                try:
                    chart.series[i].format.fill.solid()
                    chart.series[i].format.fill.fore_color.rgb = get_rgb_color(series.color)
                except:
                    pass
    
    # Show data labels if specified
    if chart_data.data_labels:
        try:
            plot = chart.plots[0]
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.show_val = True
        except:
            pass


def _create_category_chart_data(chart_data: ChartData) -> CategoryChartData:
    """Create CategoryChartData for bar, column, line, pie charts"""
    data = CategoryChartData()
    data.categories = chart_data.categories
    
    for series in chart_data.series:
        data.add_series(series.name, series.values)
    
    return data


def _create_xy_chart_data(chart_data: ChartData) -> Union[XyChartData, BubbleChartData]:
    """Create XyChartData or BubbleChartData for scatter/bubble charts"""
    if chart_data.chart_type == ChartType.BUBBLE:
        data = BubbleChartData()
        for series in chart_data.series:
            s = data.add_series(series.name)
            if series.x_values and series.sizes:
                for x, y, size in zip(series.x_values, series.values, series.sizes):
                    s.add_data_point(x, y, size)
            else:
                for i, y in enumerate(series.values):
                    x = series.x_values[i] if series.x_values else i
                    size = series.sizes[i] if series.sizes else 10
                    s.add_data_point(x, y, size)
    else:
        data = XyChartData()
        for series in chart_data.series:
            s = data.add_series(series.name)
            if series.x_values:
                for x, y in zip(series.x_values, series.values):
                    s.add_data_point(x, y)
            else:
                for i, y in enumerate(series.values):
                    s.add_data_point(i, y)
    
    return data


def _apply_chart_colors(chart, colors: List[str]):
    """Apply a custom color palette to chart series"""
    for i, color in enumerate(colors):
        if i < len(chart.series):
            try:
                chart.series[i].format.fill.solid()
                chart.series[i].format.fill.fore_color.rgb = get_rgb_color(color)
            except:
                pass


def _add_image(slide, image_elem: ImageElement):
    """Add an image to the slide"""
    # Save image temporarily
    image_stream = io.BytesIO(image_elem.image_data)
    
    shape = slide.shapes.add_picture(
        image_stream,
        Inches(image_elem.x),
        Inches(image_elem.y),
        Inches(image_elem.width),
        Inches(image_elem.height)
    )
    
    # Set rotation if specified
    if image_elem.rotation:
        shape.rotation = image_elem.rotation
    
    # Set alt text for accessibility
    if image_elem.alt_text:
        shape._element.attrib['descr'] = image_elem.alt_text


def _get_auto_shape_type(shape_type: str) -> MSO_AUTO_SHAPE_TYPE:
    """Convert shape type string to MSO_AUTO_SHAPE_TYPE"""
    shape_map = {
        'rectangle': MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        'rounded_rectangle': MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        'oval': MSO_AUTO_SHAPE_TYPE.OVAL,
        'circle': MSO_AUTO_SHAPE_TYPE.OVAL,
        'triangle': MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
        'right_triangle': MSO_AUTO_SHAPE_TYPE.RIGHT_TRIANGLE,
        'diamond': MSO_AUTO_SHAPE_TYPE.DIAMOND,
        'pentagon': MSO_AUTO_SHAPE_TYPE.PENTAGON,
        'hexagon': MSO_AUTO_SHAPE_TYPE.HEXAGON,
        'octagon': MSO_AUTO_SHAPE_TYPE.OCTAGON,
        'star': MSO_AUTO_SHAPE_TYPE.STAR_5_POINT,
        'star_4': MSO_AUTO_SHAPE_TYPE.STAR_4_POINT,
        'star_5': MSO_AUTO_SHAPE_TYPE.STAR_5_POINT,
        'star_6': MSO_AUTO_SHAPE_TYPE.STAR_6_POINT,
        'arrow_right': MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
        'arrow_left': MSO_AUTO_SHAPE_TYPE.LEFT_ARROW,
        'arrow_up': MSO_AUTO_SHAPE_TYPE.UP_ARROW,
        'arrow_down': MSO_AUTO_SHAPE_TYPE.DOWN_ARROW,
        'arrow_left_right': MSO_AUTO_SHAPE_TYPE.LEFT_RIGHT_ARROW,
        'arrow_up_down': MSO_AUTO_SHAPE_TYPE.UP_DOWN_ARROW,
        'chevron': MSO_AUTO_SHAPE_TYPE.CHEVRON,
        'callout': MSO_AUTO_SHAPE_TYPE.RECTANGULAR_CALLOUT,
        'cloud': MSO_AUTO_SHAPE_TYPE.CLOUD,
        'heart': MSO_AUTO_SHAPE_TYPE.HEART,
        'lightning': MSO_AUTO_SHAPE_TYPE.LIGHTNING_BOLT,
        'sun': MSO_AUTO_SHAPE_TYPE.SUN,
        'moon': MSO_AUTO_SHAPE_TYPE.MOON,
        'smiley': MSO_AUTO_SHAPE_TYPE.SMILEY_FACE,
        'cross': MSO_AUTO_SHAPE_TYPE.CROSS,
        'plus': MSO_AUTO_SHAPE_TYPE.CROSS,
        'block_arc': MSO_AUTO_SHAPE_TYPE.BLOCK_ARC,
        'donut': MSO_AUTO_SHAPE_TYPE.DONUT,
        'no_symbol': MSO_AUTO_SHAPE_TYPE.NO_SYMBOL,
        'trapezoid': MSO_AUTO_SHAPE_TYPE.TRAPEZOID,
        'parallelogram': MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM,
        'cube': MSO_AUTO_SHAPE_TYPE.CUBE,
        'can': MSO_AUTO_SHAPE_TYPE.CAN,
        'folded_corner': MSO_AUTO_SHAPE_TYPE.FOLDED_CORNER,
        'bevel': MSO_AUTO_SHAPE_TYPE.BEVEL,
        'frame': MSO_AUTO_SHAPE_TYPE.FRAME,
        'half_frame': MSO_AUTO_SHAPE_TYPE.HALF_FRAME,
        'plaque': MSO_AUTO_SHAPE_TYPE.PLAQUE,
        'flowchart_process': MSO_AUTO_SHAPE_TYPE.FLOWCHART_PROCESS,
        'flowchart_decision': MSO_AUTO_SHAPE_TYPE.FLOWCHART_DECISION,
        'flowchart_data': MSO_AUTO_SHAPE_TYPE.FLOWCHART_DATA,
        'flowchart_document': MSO_AUTO_SHAPE_TYPE.FLOWCHART_DOCUMENT,
        'flowchart_terminator': MSO_AUTO_SHAPE_TYPE.FLOWCHART_TERMINATOR,
    }
    return shape_map.get(shape_type.lower(), MSO_AUTO_SHAPE_TYPE.RECTANGLE)


def _add_shape(slide, shape_elem: ShapeElement):
    """Add a shape to the slide"""
    auto_shape_type = _get_auto_shape_type(shape_elem.shape_type)
    
    shape = slide.shapes.add_shape(
        auto_shape_type,
        Inches(shape_elem.x),
        Inches(shape_elem.y),
        Inches(shape_elem.width),
        Inches(shape_elem.height)
    )
    
    # Set fill color
    if shape_elem.fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = get_rgb_color(shape_elem.fill_color)
    else:
        shape.fill.background()
    
    # Set border
    if shape_elem.border_color:
        shape.line.color.rgb = get_rgb_color(shape_elem.border_color)
        shape.line.width = Pt(shape_elem.border_width)
    else:
        shape.line.fill.background()
    
    # Set rotation
    if shape_elem.rotation:
        shape.rotation = shape_elem.rotation
    
    # Add text if specified
    if shape_elem.text:
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = shape_elem.text
        p.alignment = PP_ALIGN.CENTER
        
        if shape_elem.text_style:
            font = p.font
            font.name = shape_elem.text_style.font_name
            font.size = Pt(shape_elem.text_style.font_size)
            font.bold = shape_elem.text_style.bold
            font.italic = shape_elem.text_style.italic
            font.color.rgb = get_rgb_color(shape_elem.text_style.color)


def _add_table(slide, table_elem: TableElement):
    """Add a table to the slide"""
    rows = len(table_elem.rows)
    cols = max(len(row) for row in table_elem.rows) if table_elem.rows else 1
    
    table_shape = slide.shapes.add_table(
        rows, cols,
        Inches(table_elem.x),
        Inches(table_elem.y),
        Inches(table_elem.width),
        Inches(table_elem.height)
    )
    
    table = table_shape.table
    
    # Fill in cells
    for row_idx, row_data in enumerate(table_elem.rows):
        for col_idx, cell_data in enumerate(row_data):
            if col_idx >= cols:
                continue
            
            cell = table.cell(row_idx, col_idx)
            
            # Set text
            text_frame = cell.text_frame
            text_frame.word_wrap = True
            p = text_frame.paragraphs[0]
            p.text = cell_data.text
            p.alignment = _get_alignment(cell_data.style.alignment)
            
            # Apply font styling
            if p.runs:
                font = p.runs[0].font
                font.name = cell_data.style.font_name
                font.size = Pt(cell_data.style.font_size)
                font.bold = cell_data.style.bold
                font.italic = cell_data.style.italic
                font.color.rgb = get_rgb_color(cell_data.style.color)
            
            # Set background color
            if cell_data.background_color:
                cell.fill.solid()
                cell.fill.fore_color.rgb = get_rgb_color(cell_data.background_color)
            elif row_idx == 0 and table_elem.header_row and table_elem.header_color:
                cell.fill.solid()
                cell.fill.fore_color.rgb = get_rgb_color(table_elem.header_color)
            elif table_elem.alternating_row_color and row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = get_rgb_color(table_elem.alternating_row_color)


# ============================================================================
# High-Level Helper Functions
# ============================================================================

def create_title_slide(
    title: str,
    subtitle: str = "",
    title_style: Optional[TextStyle] = None,
    subtitle_style: Optional[TextStyle] = None,
    background_color: str = "#FFFFFF"
) -> SlideContent:
    """Create a simple title slide"""
    if title_style is None:
        title_style = TextStyle(
            font_name="Calibri",
            font_size=44,
            bold=True,
            color="#000000",
            alignment="center"
        )
    
    if subtitle_style is None:
        subtitle_style = TextStyle(
            font_name="Calibri",
            font_size=24,
            color="#666666",
            alignment="center"
        )
    
    title_box = TextBox(
        paragraphs=[
            Paragraph(
                runs=[TextRun(title, title_style)],
                alignment="center"
            )
        ],
        x=0.5, y=2.5, width=12.333, height=1.5
    )
    
    subtitle_box = TextBox(
        paragraphs=[
            Paragraph(
                runs=[TextRun(subtitle, subtitle_style)],
                alignment="center"
            )
        ],
        x=0.5, y=4.0, width=12.333, height=1.0
    ) if subtitle else None
    
    layout = SlideLayout(background_color=background_color)
    
    text_boxes = [title_box]
    if subtitle_box:
        text_boxes.append(subtitle_box)
    
    return SlideContent(layout=layout, text_boxes=text_boxes)


def create_content_slide(
    title: str,
    bullet_points: List[str],
    title_style: Optional[TextStyle] = None,
    bullet_style: Optional[TextStyle] = None,
    background_color: str = "#FFFFFF"
) -> SlideContent:
    """Create a slide with title and bullet points"""
    if title_style is None:
        title_style = TextStyle(
            font_name="Calibri",
            font_size=32,
            bold=True,
            color="#000000"
        )
    
    if bullet_style is None:
        bullet_style = TextStyle(
            font_name="Calibri",
            font_size=20,
            color="#333333"
        )
    
    title_box = TextBox(
        paragraphs=[
            Paragraph(
                runs=[TextRun(title, title_style)],
                alignment="left"
            )
        ],
        x=0.5, y=0.5, width=12.333, height=1.0
    )
    
    bullet_paragraphs = [
        Paragraph(
            runs=[TextRun(f"• {point}", bullet_style)],
            alignment="left",
            bullet=True
        )
        for point in bullet_points
    ]
    
    content_box = TextBox(
        paragraphs=bullet_paragraphs,
        x=0.5, y=1.7, width=12.333, height=5.0
    )
    
    layout = SlideLayout(background_color=background_color)
    
    return SlideContent(layout=layout, text_boxes=[title_box, content_box])


def create_chart_slide(
    title: str,
    chart_data: ChartData,
    title_style: Optional[TextStyle] = None,
    background_color: str = "#FFFFFF"
) -> SlideContent:
    """Create a slide with a chart"""
    if title_style is None:
        title_style = TextStyle(
            font_name="Calibri",
            font_size=32,
            bold=True,
            color="#000000"
        )
    
    title_box = TextBox(
        paragraphs=[
            Paragraph(
                runs=[TextRun(title, title_style)],
                alignment="left"
            )
        ],
        x=0.5, y=0.5, width=12.333, height=1.0
    )
    
    chart_elem = ChartElement(
        data=chart_data,
        x=0.5, y=1.7, width=12.333, height=5.3
    )
    
    layout = SlideLayout(background_color=background_color)
    
    return SlideContent(layout=layout, text_boxes=[title_box], charts=[chart_elem])


def create_image_slide(
    title: str,
    image_data: bytes,
    title_style: Optional[TextStyle] = None,
    background_color: str = "#FFFFFF"
) -> SlideContent:
    """Create a slide with a full-size image"""
    if title_style is None:
        title_style = TextStyle(
            font_name="Calibri",
            font_size=32,
            bold=True,
            color="#000000"
        )
    
    title_box = TextBox(
        paragraphs=[
            Paragraph(
                runs=[TextRun(title, title_style)],
                alignment="left"
            )
        ],
        x=0.5, y=0.5, width=12.333, height=1.0
    )
    
    image_elem = ImageElement(
        image_data=image_data,
        x=0.5, y=1.7, width=12.333, height=5.3
    )
    
    layout = SlideLayout(background_color=background_color)
    
    return SlideContent(layout=layout, text_boxes=[title_box], images=[image_elem])


def image_to_bytes(image: Image.Image, format: str = "PNG") -> bytes:
    """Convert PIL Image to bytes"""
    buffer = io.BytesIO()
    image.save(buffer, format=format)
    return buffer.getvalue()


def create_slide_from_image(
    image: Image.Image,
    slide_width: float = 13.333,
    slide_height: float = 7.5
) -> SlideContent:
    """
    Create a slide with the image as background (fallback for complex slides).
    Used when AI extraction fails and we need to preserve the visual exactly.
    """
    image_bytes = image_to_bytes(image)
    
    # Calculate image dimensions to fit slide while maintaining aspect ratio
    img_aspect = image.width / image.height
    slide_aspect = slide_width / slide_height
    
    if img_aspect > slide_aspect:
        # Image is wider - fit to width
        width = slide_width
        height = slide_width / img_aspect
        x = 0
        y = (slide_height - height) / 2
    else:
        # Image is taller - fit to height
        height = slide_height
        width = slide_height * img_aspect
        x = (slide_width - width) / 2
        y = 0
    
    image_elem = ImageElement(
        image_data=image_bytes,
        x=x, y=y,
        width=width, height=height
    )
    
    layout = SlideLayout(width=slide_width, height=slide_height)
    
    return SlideContent(layout=layout, images=[image_elem])
